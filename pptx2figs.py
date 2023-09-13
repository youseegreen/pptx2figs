from genericpath import isfile
import os
import shutil
import sys
import getopt
try:
    from pptx import Presentation
except:
    print("you have no python-pptx!")
try:
    import win32com.client
except:
    print("you have no pywin32!")


usage = '''pptx2figs.py usage:

python pptx2figs.py <options>
        --input=<str> / -i : 対象のpptxファイルネーム
        --start=<number> / -s : 探索開始のスライド番号 (default : 0)
        --end=<number> / -e : 探索終了のスライド番号 (default : None)
'''

# 保存先のフォルダ、上書きしてしまうので注意
f_pptx = "./pptxs/"
f_pdf = "./pdfs/"
f_png = "./pngs/"


''' 入力pptxファイルから対象図のみのpptxファイルを作成する '''
def make_1fig_pptx(prs, pptx_filename, t_page, t_shape_id):
    # shapeオブジェクトをコピーして代入が上手く動かないので、prsを削る方向で進める
    # dst = copy.deepcopy(prs)だと上手く動かなかったので、一旦別ファイルとして保存する
    prs.save(pptx_filename)
    dst = Presentation(pptx_filename) 

    # 情報取得
    t_shape = dst.slides[t_page].shapes[t_shape_id]
    bias_x = t_shape.left
    bias_y = t_shape.top
    width = t_shape.width
    height = t_shape.height

    # スライドサイズ変更
    dst.slide_width = width
    dst.slide_height = height

    # t_page以外は削除
    xml_slides = dst.slides._sldIdLst
    slides = list(xml_slides)
    for page in range(len(dst.slides) - 1, -1, -1):
        if page != t_page:
            rId = dst.slides._sldIdLst[page].rId
            dst.part.drop_rel(rId)
            xml_slides.remove(slides[page])

    # (枠+pdf_name)shapeを削除する
    shapes = dst.slides[0].shapes   # target slideは0ページになっている
    shapes.element.remove(shapes[t_shape_id].element)  

    # 全てのオブジェクトの原点を変更
    for shape in shapes:
        shape.left -= bias_x
        shape.top -= bias_y

    # pptx_filenameに上書き保存
    dst.save(pptx_filename)    


''' pptxファイルからpdfファイルを作成する '''
def make_1fig_pdf(pptx_filename, pdf_filename):

    # 参考：https://python-work.com/pptx-to-pdf/
    # 絶対パスで設定しないと動かない
    pptx_filename_abs = os.path.join(os.path.abspath("./"), pptx_filename)
    pdf_filename_abs = os.path.join(os.path.abspath("./"), pdf_filename)

    # Powerpointファイルを開きPDF形式で保存
    application = win32com.client.Dispatch('Powerpoint.Application')
    presentation = application.Presentations.Open(pptx_filename_abs)
    presentation.SaveAs(pdf_filename_abs, 32) 
    
    # アプリケーション終了処理
    presentation.close()
    application.quit()
    presentation = None
    application = None


''' pptxファイルからpngファイルを作成する '''
def make_1fig_png(pptx_filename, png_filename):

    # # 絶対パスで設定しないと動かない
    pptx_filename_abs = os.path.join(os.path.abspath("./"), pptx_filename)
    png_filename_abs = os.path.join(os.path.abspath("./"), png_filename)

    # Powerpointファイルを開きPDF形式で保存
    application = win32com.client.Dispatch('Powerpoint.Application')
    presentation = application.Presentations.Open(pptx_filename_abs)
    presentation.Export(png_filename_abs, FilterName="png")

    # HOGEHOGE.pngを保存したい場合、HOGEHOGE/スライド1.PNGになるので、修正する
    shutil.move(png_filename[:-4] + "/スライド1.PNG", png_filename)
    os.rmdir(png_filename[:-4])

    # アプリケーション終了処理
    presentation.close()
    application.quit()
    presentation = None
    application = None


''' 入力pptxファイルから図にしたいとこを全て抽出してpptx, pdf, pngを作成する '''
def make_figs(target_filename, start_page = 0, end_page = None):

    # 保存先のフォルダを作っておく
    if not os.path.exists(f_pptx):
        os.makedirs(f_pptx)
    if not os.path.exists(f_pdf):
        os.makedirs(f_pdf)
    if not os.path.exists(f_png):
        os.makedirs(f_png)

    # target_filename (.pptx)を開く
    if target_filename[-5:] != ".pptx":
        print(f"{target_filename} is not a pptx file!")
        return 
    if not os.path.isfile(target_filename):
        print(f"{target_filename} is not found!")
        return 
    prs = Presentation(target_filename)

    # textが*.pdfのshapeがあれば、pptx, pdfを作成する
    print("[Targets]")
    for page, slide in enumerate(prs.slides):

        if page < start_page or (end_page != None and page > end_page):
            continue

        for shape_id, shape in enumerate(slide.shapes):
            try:
                if '.pdf' in shape.text:  # 四角図形の場合も加えたい
                    print(f" - slide {page} : {shape.text}")
                    pptx_filename = f_pptx + shape.text[:-4] + ".pptx"
                    pdf_filename = f_pdf + shape.text
                    png_filename = f_png + shape.text[:-4] + ".png"

                    # 対象figureだけのpptxを作成
                    make_1fig_pptx(prs, pptx_filename, page, shape_id)
                    # 対象figureだけのpptxからpdfを作成
                    make_1fig_pdf(pptx_filename, pdf_filename)
                    # 対象figureだけのpptxからpngを作成
                    make_1fig_png(pptx_filename, png_filename)

            except:
                pass


if __name__ == "__main__":

    start_page = 0
    end_page = None
    input_filename = None

    argv = sys.argv[1:]
    
    try:
        opts, args = getopt.getopt(argv, 'h:i:s:e', ['help', 'input=', 'start=', 'end='])
    except getopt.GetoptError:
        print(usage)
        sys.exit()
    for opt, arg in opts:
        try:
            if opt in ('-h', '--help'):
                print(usage)
                sys.exit()
            elif opt in ('-i', '--input'):
                input_filename = arg
            elif opt in ('-s', '--start'):
                start_page = int(arg)
            elif opt in ('-e', '--end'):
                end_page = int(arg)
        except Exception:
            print('Error parsing argument: %s' % opt)
            print(usage)
            sys.exit(2)    
    if input_filename is not None:   
        make_figs(input_filename, start_page, end_page)
    else:
        f_input = "./input_pptxs/"
        files = [f_input + f for f in os.listdir(f_input) if '.pptx' in f]
        for file in files:
            make_figs(file, start_page, end_page)